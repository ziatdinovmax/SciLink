"""Generate a 5-slide overview deck of SciLink's simulation agent stack.

Audience: internal team / collaborators. Sets up the validation panel results
that are currently running on Deception.

Usage:
    python build_overview_deck.py
    # → scilink_sim_agents_overview.pptx
"""
from datetime import date

from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor


# ── theme ─────────────────────────────────────────────────────────
DARK = RGBColor(0x1F, 0x2D, 0x3D)        # title / heavy text
MID = RGBColor(0x4A, 0x5B, 0x6E)         # body
ACCENT = RGBColor(0x2E, 0x86, 0xAB)      # PNNL-ish blue
MUTED = RGBColor(0x88, 0x90, 0x9C)       # placeholder / tbd
FILL_HEAD = RGBColor(0x2E, 0x86, 0xAB)
FILL_ALT = RGBColor(0xF2, 0xF5, 0xF8)

# Schematic node fills — pulled from the same family but distinguishable
INPUT_FILL = RGBColor(0xFE, 0xEA, 0xC2)      # warm yellow — external input
ORCH_FILL = RGBColor(0xD7, 0xE8, 0xF3)       # light accent blue — orchestrator core
ROUTER_FILL = RGBColor(0xBC, 0xD9, 0xEB)     # mid accent blue — router
AGENT_FILL = RGBColor(0xFF, 0xFF, 0xFF)      # white — peer scale agents
OUTPUT_FILL = RGBColor(0xEC, 0xEE, 0xF1)     # cool gray — outputs
SKILL_FILL = RGBColor(0xE3, 0xD4, 0xE3)      # soft mauve — skill subsystem
SKILL_LINE = RGBColor(0x7C, 0x52, 0x7C)      # darker mauve border


def set_widescreen(prs: Presentation) -> None:
    prs.slide_width = Emu(12192000)   # 13.333"
    prs.slide_height = Emu(6858000)   # 7.5"


def add_title(slide, text: str, *, top=Inches(0.35), size=32) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.9))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = DARK


def add_subtitle(slide, text: str, *, top=Inches(1.15), size=16,
                 color=MID) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.italic = True
    run.font.color.rgb = color


def add_bullets(slide, bullets, *, left=Inches(0.5), top=Inches(1.6),
                width=Inches(12.3), height=Inches(5.5), size=16) -> None:
    """bullets: list of (level, text) or plain str (level 0)."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        prefix = "•  " if level == 0 else "–  "
        p.text = prefix + text
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.size = Pt(size - 2 * level)
        run.font.color.rgb = DARK if level == 0 else MID


def add_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text.strip()


def add_node(slide, x, y, w, h, label, *, sublabel=None,
             sublabel_color=None, sublabel_size=None,
             fill=FILL_ALT, border=ACCENT, font_color=DARK,
             font_size=12, bold=False, italic=False,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Inches in / inches in. Returns the shape so callers can grab edges."""
    sh = slide.shapes.add_shape(
        shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    for i, line in enumerate(label.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.text = line
        if not p.runs:
            continue
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = font_color
    if sublabel:
        for line in sublabel.split("\n"):
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.text = line
            if not p2.runs:
                continue
            r2 = p2.runs[0]
            r2.font.size = Pt(sublabel_size if sublabel_size else font_size - 2)
            r2.font.italic = True
            r2.font.color.rgb = sublabel_color if sublabel_color else MID
    return sh


def add_arrow(slide, x1, y1, x2, y2, *, color=DARK, width_pt=1.5,
              dashed=False, arrow=True):
    """Inches. Straight connector; arrowhead at (x2,y2) when arrow=True.

    Pass arrow=False for plain line segments (e.g. middle pieces of a
    multi-segment elbow loop, where only the final segment carries the head).
    """
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    ln = conn.line._get_or_add_ln()
    if arrow:
        tail = ln.find(qn('a:tailEnd'))
        if tail is None:
            tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
        tail.set('w', 'med')
        tail.set('len', 'med')
    if dashed:
        prst = ln.find(qn('a:prstDash'))
        if prst is None:
            prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', 'dash')
    return conn


def add_caption(slide, x, y, w, h, text, *, size=11, color=MID,
                italic=True, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.text = line
        if not p.runs:
            continue
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def style_table(table, *, header_fill=FILL_HEAD, alt_fill=FILL_ALT,
                head_color=RGBColor(0xFF, 0xFF, 0xFF), body_color=DARK,
                font_size=12) -> None:
    nrows = len(table.rows)
    ncols = len(table.columns)
    for r in range(nrows):
        for c in range(ncols):
            cell = table.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = head_color if r == 0 else body_color
                    if r == 0:
                        run.font.bold = True


def fill_table(table, rows) -> None:
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)


# ── slide builders ────────────────────────────────────────────────

def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Accent bar
    from pptx.shapes.autoshape import Shape
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.1), Inches(13.333), Inches(0.08),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Title
    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    p.text = "SciLink Simulation Agents"
    run = p.runs[0]
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = DARK

    box = slide.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(12), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = "Capability overview & validation in progress"
    run = p.runs[0]
    run.font.size = Pt(22)
    run.font.italic = True
    run.font.color.rgb = MID

    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(12), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = f"Sarah Allec  ·  PNNL  ·  {date.today().isoformat()}"
    run = p.runs[0]
    run.font.size = Pt(13)
    run.font.color.rgb = MUTED

    add_notes(slide, """
    Overview deck for internal use. Sets up the validation results that are
    running on Deception today (the a/b/c lattice-constant panel). Covers
    architecture, the two pipelines, the autonomy spectrum, and a results
    scaffold.
    """)


def slide_stack(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Inside the simulate mode — agent architecture")
    add_subtitle(slide,
                 "One of three top-level modes (analyze · plan · simulate). "
                 "This is what happens after a user goal enters the simulate side.")

    # ── schematic layout (inches) ──
    # Skills (top, wide bar spanning orchestrator + agents region)
    add_node(slide, 2.1, 1.85, 7.9, 0.7,
             "Skill subsystem  ·  scilink/skills/  +  ~/.scilink/graduated_skills/",
             sublabel="domain-keyed bundles  →  orchestrator routing  +  "
                      "per-agent engine support",
             fill=SKILL_FILL, border=SKILL_LINE, bold=True, font_size=13,
             sublabel_color=SKILL_LINE, sublabel_size=11)

    # User goal (far left)
    add_node(slide, 0.4, 4.0, 1.4, 0.9,
             "User goal",
             sublabel="natural language",
             fill=INPUT_FILL, border=RGBColor(0xC8, 0xA8, 0x3C),
             bold=True, font_size=12)

    # Orchestrator (center)
    add_node(slide, 2.1, 3.85, 3.1, 1.2,
             "SimulationOrchestratorAgent",
             sublabel="run_chat()  ·  run_task()",
             fill=ORCH_FILL, border=ACCENT, bold=True, font_size=14)

    # Router
    add_node(slide, 5.55, 4.05, 1.5, 0.85,
             "Router",
             sublabel="(scale, engine)",
             fill=ROUTER_FILL, border=ACCENT, bold=True, font_size=12)

    # Scale agents (stacked) — each shows its supported engines in mauve
    # to visually link "what each agent can drive" back to the skill subsystem.
    agent_eng_size = 10
    add_node(slide, 7.4, 3.05, 2.6, 0.85,
             "PeriodicDFTAgent",
             sublabel="VASP · QE · ABINIT · CP2K",
             fill=AGENT_FILL, border=ACCENT, bold=True, font_size=12,
             sublabel_color=SKILL_LINE, sublabel_size=agent_eng_size)
    add_node(slide, 7.4, 4.05, 2.6, 0.85,
             "MDSimulationAgent",
             sublabel="LAMMPS · GROMACS · OpenMM · AMBER",
             fill=AGENT_FILL, border=ACCENT, bold=True, font_size=12,
             sublabel_color=SKILL_LINE, sublabel_size=agent_eng_size)
    add_node(slide, 7.4, 5.05, 2.6, 0.85,
             "MLIPAgent",
             sublabel="MACE · NequIP · CHGNet · DeePMD",
             fill=AGENT_FILL, border=ACCENT, bold=True, font_size=12,
             sublabel_color=SKILL_LINE, sublabel_size=agent_eng_size)

    # Outputs (far right) — concrete examples; engine-specific
    add_node(slide, 10.25, 3.125, 2.6, 0.7,
             "POSCAR · INCAR · KPOINTS",
             sublabel="(VASP example)",
             fill=OUTPUT_FILL, border=MUTED, font_size=11)
    add_node(slide, 10.25, 4.125, 2.6, 0.7,
             "MD script · trajectory",
             sublabel="(LAMMPS example)",
             fill=OUTPUT_FILL, border=MUTED, font_size=11)
    add_node(slide, 10.25, 5.125, 2.6, 0.7,
             "DeployedPotential",
             sublabel="(engine-neutral descriptor)",
             fill=OUTPUT_FILL, border=MUTED, font_size=11)

    # ── arrows ──
    # skill → orchestrator (dashed: "context", not data flow)
    add_arrow(slide, 3.65, 2.55, 3.65, 3.85, color=SKILL_LINE, dashed=True)
    # skill → agents column (dashed: feeds each agent's engine support)
    add_arrow(slide, 8.7, 2.55, 8.7, 3.05, color=SKILL_LINE, dashed=True)
    # user → orchestrator
    add_arrow(slide, 1.8, 4.45, 2.1, 4.45)
    # orchestrator → router
    add_arrow(slide, 5.2, 4.45, 5.55, 4.475)
    # router → 3 agents (fan out)
    add_arrow(slide, 7.05, 4.475, 7.4, 3.475)
    add_arrow(slide, 7.05, 4.475, 7.4, 4.475)
    add_arrow(slide, 7.05, 4.475, 7.4, 5.475)
    # agents → outputs
    add_arrow(slide, 10.0, 3.475, 10.25, 3.475)
    add_arrow(slide, 10.0, 4.475, 10.25, 4.475)
    add_arrow(slide, 10.0, 5.475, 10.25, 5.475)

    # ── caption / legend ──
    add_caption(slide, 0.5, 6.2, 12.3, 0.4,
                "Solid arrows  =  data flow.    "
                "Dashed arrows  =  skill-bundle context — feeding the "
                "router's scale decision  and  each scale agent's engine set.")
    add_caption(slide, 0.5, 6.55, 12.3, 0.4,
                "Scale agents are software-agnostic — each routes to one of "
                "several engines; adding an engine is one skill bundle, "
                "no agent-code edit.",
                size=11)

    add_notes(slide, """
    The three-mode commitment is in CLAUDE.md. Note the doc still says
    SimulationOrchestratorAgent is "not yet built" — it actually exists and
    is what this deck describes.

    The run_task contract returns {status, summary, files_produced,
    key_findings, suggested_followups, structures, warnings}. That's the
    surface a future meta-agent will consume.

    Skill graduation: agents record session-scoped observations into a
    KnowledgeStore (auto-UUIDs), and successful patterns can be promoted
    to persistent .md files that load on next start.
    """)


def slide_pipelines(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Self-refinement  —  run, branch, iterate")
    add_subtitle(slide,
                 "Generate runs once.  The loop is around Run.  "
                 "Engine error → debugger fires.  Phase success → quality check fires.  "
                 "Iterate until all phases pass.")

    # ── Top row: linear forward flow ──
    top_y, top_h = 1.95, 0.80
    mid_y = top_y + top_h / 2

    add_node(slide, 0.7, top_y, 2.2, top_h,
             "Generate",
             sublabel="structure  +  inputs\n(with internal refinement)",
             fill=ORCH_FILL, border=ACCENT, bold=True, font_size=14,
             sublabel_size=9)

    add_node(slide, 3.5, top_y, 3.2, top_h,
             "Run   (multi-phase)",
             sublabel="Optim  →  Equilib  →  Production",
             fill=ORCH_FILL, border=ACCENT, bold=True, font_size=14,
             sublabel_size=10)

    add_node(slide, 7.3, top_y, 2.7, top_h,
             "Quality check",
             sublabel="fires after every phase",
             fill=ORCH_FILL, border=ACCENT, bold=True, font_size=14,
             sublabel_size=10)

    add_node(slide, 10.7, top_y, 1.6, top_h,
             "Done",
             sublabel="(all phases passed)",
             fill=OUTPUT_FILL, border=MUTED, bold=True, font_size=14,
             sublabel_size=9)

    # Forward arrows along the top row
    add_arrow(slide, 2.9, mid_y, 3.5, mid_y)            # Generate → Run
    add_arrow(slide, 6.7, mid_y, 7.3, mid_y)            # Run → Quality
    add_caption(slide, 6.7, mid_y - 0.32, 0.6, 0.22,
                "no error", size=9, italic=True, color=MID,
                align=PP_ALIGN.CENTER)
    add_arrow(slide, 10.0, mid_y, 10.7, mid_y)          # Quality → Done
    add_caption(slide, 10.0, mid_y - 0.32, 0.7, 0.22,
                "all pass", size=9, italic=True, color=MID,
                align=PP_ALIGN.CENTER)

    # ── Debugger box below Run (the error branch) ──
    dbg_y, dbg_h = 3.55, 0.70
    add_node(slide, 3.9, dbg_y, 2.4, dbg_h,
             "Debugger",
             sublabel="VaspUpdater / LAMMPSUpdater\n"
                      "parses engine error  →  corrected params",
             fill=ORCH_FILL, border=ACCENT, bold=True, font_size=13,
             sublabel_size=8)

    # Run → Debugger  (solid, on engine error)
    add_arrow(slide, 5.8, top_y + top_h, 5.8, dbg_y)
    add_caption(slide, 5.9, 2.92, 1.3, 0.22,
                "engine error", size=9, italic=True, color=MID,
                align=PP_ALIGN.LEFT)

    # Debugger → Run  (dashed mauve, corrected params)
    add_arrow(slide, 4.4, dbg_y, 4.4, top_y + top_h,
              color=SKILL_LINE, dashed=True, arrow=True)
    add_caption(slide, 2.95, 2.92, 1.4, 0.22,
                "corrected params", size=9, italic=True, color=SKILL_LINE,
                align=PP_ALIGN.RIGHT)

    # ── Quality → Run  (dashed mauve elbow, refine if questionable) ──
    # Routes ABOVE the Debugger box so it doesn't intersect it (only
    # crosses the two vertical Run↔Debugger arrows — different style/color
    # makes the crossing readable).
    qx = 8.65               # Quality bottom-center x
    loop_y = 3.25           # horizontal segment y (above Debugger top at 3.55)
    refine_x = 3.7          # entry x into Run bottom-left
    add_arrow(slide, qx, top_y + top_h, qx, loop_y,
              color=SKILL_LINE, dashed=True, arrow=False)
    add_arrow(slide, qx, loop_y, refine_x, loop_y,
              color=SKILL_LINE, dashed=True, arrow=False)
    add_arrow(slide, refine_x, loop_y, refine_x, top_y + top_h,
              color=SKILL_LINE, dashed=True, arrow=True)
    add_caption(slide, 6.5, loop_y + 0.05, 3.5, 0.22,
                "refine if questionable",
                size=9, italic=True, color=SKILL_LINE,
                align=PP_ALIGN.CENTER)

    # ── Per-scale instance rows (4 columns: Generate, Run, Debugger, Quality) ──
    COL_X_PS = [1.20, 4.10, 7.00, 9.90]
    COL_W_PS = 2.80

    # column header strip — small labels above the per-scale grid
    header_y = 4.55
    for x, label in zip(COL_X_PS,
                        ["Generate", "Run", "Debugger", "Quality"]):
        add_caption(slide, x, header_y, COL_W_PS, 0.25,
                    label, size=10, italic=True, color=MID,
                    align=PP_ALIGN.CENTER)

    SCALE_ROWS = [
        ("DFT",  4.85, [
            "StructureGenerator\n+ VaspInputAgent",
            "VASP  (multi-block ionic relax)",
            "VaspUpdater\nparses error log → fix params",
            "VaspQualityAgent\nper ionic block",
        ]),
        ("MD",   5.45, [
            "plan_simulation\n+ ForceFieldAgent + Packmol",
            "LAMMPS  (optim → equilib → prod)",
            "LAMMPSUpdater\nparses engine error → fix script",
            "LAMMPSAnalysisAgent\nper phase",
        ]),
        ("MLIP", 6.05, [
            "MLIPAgent.deploy_pretrained",
            "MD via DeployedPotential  (phased)",
            "per-step UQ\nflags diverging trajectories",
            "evaluate_simulation_quality\n+ UQ per phase",
        ]),
    ]
    row_h = 0.50

    # row labels at the left
    for label, y, _ in SCALE_ROWS:
        add_caption(slide, 0.4, y + row_h / 2 - 0.15, 0.70, 0.3,
                    label, size=12, italic=False, color=DARK,
                    align=PP_ALIGN.RIGHT)

    # cells
    for _, y, cells in SCALE_ROWS:
        for x, txt in zip(COL_X_PS, cells):
            add_node(slide, x, y, COL_W_PS, row_h, txt,
                     fill=FILL_ALT, border=MUTED,
                     font_size=9, bold=False, italic=False,
                     font_color=DARK)

    # ── Bottom captions ──
    add_caption(slide, 0.5, 6.70, 12.3, 0.30,
                "Same loop across DFT, classical MD, and MLIP-driven MD — "
                "only the agents in each slot differ.",
                size=10)

    add_notes(slide, """
    User correction (v3 of this slide):  the iteration is around RUN,
    not around Generate.  Generate runs ONCE at the start.

    After each Run, branch on outcome:
      - Engine error  →  Debugger parses the error log and proposes
                         corrected input params  →  back to Run.
      - Phase success →  Quality check assesses the physics  →  pass
                         continues to the next phase (or DONE if all
                         phases are complete); questionable triggers
                         a refine back to Run with adjusted params.

    Both feedback paths terminate at Run, not at Generate.  That's why
    the loops in the schematic are tight around Run rather than long
    arcs back to the leftmost stage.

    Pre-run validators (StructureValidatorAgent, IncarValidatorAgent)
    live INSIDE the Generate stage — they refine the structure / inputs
    as part of generation.  They are not a separate pre-run stage.

    Visual note:  the Quality → Run loop crosses the two vertical
    Run↔Debugger arrows at y=3.25.  Different colors (dark solid vs
    mauve dashed) keep the crossing readable.
    """)


def slide_autonomy(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Autonomy spectrum — what the benchmark covers")
    add_subtitle(slide,
                 "Two orthogonal axes the agent reasons across.  "
                 "Today's router test exercises decisions before the calc; "
                 "cluster harnesses (test_dft, test_mlip) test decisions within it.")

    # ── Grid geometry ──
    col_x = [3.1, 5.65, 8.2]       # left edges of three columns
    row_y = [2.55, 3.65, 4.75]     # top edges (top row = autonomous)
    cell_w, cell_h = 2.45, 1.00
    col_labels = ["a_forced",  "b_agent_select",  "c_bare_goal"]
    row_labels = ["autonomous", "supervised",     "co-pilot"]

    # ── X-axis label (top) ──
    add_caption(slide, col_x[0], 1.7, col_x[-1] + cell_w - col_x[0], 0.3,
                "Reasoning autonomy  —  how much the agent decides about the "
                "calc  →",
                size=12, italic=True, color=MID)

    # ── Column headers (just above row 1) ──
    for x, label in zip(col_x, col_labels):
        add_caption(slide, x, 2.10, cell_w, 0.35,
                    label, size=13, italic=False, color=DARK,
                    align=PP_ALIGN.CENTER)

    # ── Y-axis label (left, vertical stack) ──
    add_caption(slide, 0.4, 2.55, 1.3, 0.35,
                "Runtime autonomy  ↑",
                size=12, italic=True, color=MID, align=PP_ALIGN.LEFT)
    add_caption(slide, 0.4, 2.85, 1.3, 0.35,
                "(approval cadence per tool call)",
                size=9, italic=True, color=MUTED, align=PP_ALIGN.LEFT)

    # ── Row labels (right-aligned, left of grid) ──
    for y, label in zip(row_y, row_labels):
        add_caption(slide, 1.85, y + cell_h / 2 - 0.18, 1.15, 0.35,
                    label, size=13, italic=False, color=DARK,
                    align=PP_ALIGN.RIGHT)

    # ── 3 × 3 grid of cells ──
    # Conceptual map only — cells are left empty; the benchmark suite
    # exercises individual cells (initially the autonomous × c_bare_goal
    # corner via test_dft and test_mlip) but we don't claim coverage
    # yet on this slide.
    for r in range(3):
        for c in range(3):
            add_node(slide, col_x[c], row_y[r], cell_w, cell_h,
                     "",
                     fill=RGBColor(0xFF, 0xFF, 0xFF), border=MUTED,
                     font_size=14, bold=False,
                     font_color=MUTED)

    # ── Bottom captions ──
    add_caption(slide, 0.5, 5.95, 12.3, 0.35,
                "Reasoning autonomy is what the agent decides about the calc "
                "(backend, params, structure). Runtime autonomy is how often "
                "the human gets to approve.",
                size=11)
    add_caption(slide, 0.5, 6.35, 12.3, 0.35,
                "The benchmark sweeps the reasoning axis; the runtime axis is "
                "set per session.  Filling more of this grid is the roadmap.",
                size=11)

    add_notes(slide, """
    Two orthogonal axes:

      * Reasoning autonomy (columns) — how much of the calculation the
        agent decides on its own.  a_forced names the backend + gives
        a guided goal; b_agent_select lets the agent pick the backend
        but still hands it physics hints; c_bare_goal gives only the
        material and the task.

      * Runtime autonomy (rows) — how often the human approves.
        co-pilot gates every tool call; supervised surfaces key
        decisions; autonomous runs end-to-end without confirmation.

    These are independent — any (column, row) pair is a valid session.
    The current benchmark suite (validation-suite branch) sweeps the
    columns at the autonomous row.  Filling the lower rows is roadmap;
    the orchestrator already supports them, the harness work is just
    to script "approve the tool call" for the gated modes.
    """)


def slide_router(prs: Presentation) -> None:
    """Slide 5/6: router decision quality."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Router  —  picks (scale, engine) from a prompt")
    add_subtitle(slide,
                 "Materials problems are scale-ambiguous:  lattice constants "
                 "are DFT, melting points need MLIP / MD, gas-phase chemistry "
                 "needs molecular DFT.  The router has to pick the right "
                 "tool from a natural-language prompt.")

    # ── Top: 3 headline number cards ──
    def _stat_card(x: float, big: str, label: str, accent=ACCENT) -> None:
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(1.7), Inches(4.0), Inches(1.0))
        bg.fill.solid()
        bg.fill.fore_color.rgb = ORCH_FILL
        bg.line.color.rgb = accent
        bg.shadow.inherit = False
        # big number
        box = slide.shapes.add_textbox(
            Inches(x), Inches(1.78), Inches(4.0), Inches(0.55))
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = big
        run = p.runs[0]
        run.font.size = Pt(32); run.font.bold = True; run.font.color.rgb = accent
        # label
        box2 = slide.shapes.add_textbox(
            Inches(x), Inches(2.30), Inches(4.0), Inches(0.4))
        p2 = box2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        p2.text = label
        run2 = p2.runs[0]
        run2.font.size = Pt(11); run2.font.italic = True
        run2.font.color.rgb = DARK

    _stat_card(0.5, "100 %", "joint accuracy   (3 models × 21 scorable prompts)")
    _stat_card(4.7, "1.00",  "variability stability   (8-shot test × 3 prompts × 2 models)",
               accent=ACCENT)
    _stat_card(8.9, "1",     "stable failure mode caught by variability",
               accent=SKILL_LINE)

    # ── Middle: Example prompts & outcomes ──
    add_caption(slide, 0.5, 2.95, 12.3, 0.28,
                "Example prompts & outcomes   (selected from 24-query suite)",
                size=12, italic=False, color=ACCENT, align=PP_ALIGN.LEFT)

    tbl = slide.shapes.add_table(
        7, 3,
        Inches(0.5), Inches(3.25), Inches(12.3), Inches(2.50)).table
    for c, w in enumerate([7.4, 2.3, 2.6]):
        tbl.columns[c].width = Inches(w)
    fill_table(tbl, [
        ["user prompt",
         "expected",
         "agent's route"],
        ['pd_07_battery_licoo2 — "Relax the LiCoO₂ cathode crystal in its layered R-3m phase."',
         "periodic_dft / vasp",
         "periodic_dft / vasp   ✓"],
        ['cmd_03_battery_electrolyte — "Build an MD simulation of 1 M LiPF₆ in ethylene carbonate."',
         "MD / any FF engine",
         "MD / lammps   ✓"],
        ['mlip_02_battery_cathode — "Run MD on a LiCoO₂ cathode cell at elevated T with a pretrained ML potential."',
         "MLIP / any",
         "MLIP / chgnet   ✓"],
        ['amb_01_li_diffusion — "Compute the diffusion coefficient of Li in LiCoO₂ at elevated temperature."',
         "MLIP / any  (ambiguous)",
         "MLIP / chgnet   ✓"],
        ['amb_03_melting_point_cu — "Predict the melting point of Cu."',
         "MLIP / MD only\n(DFT can\'t sample liquid)",
         "opus-4-6 → MLIP/mace  ✓ 8/8\nsonnet → periodic_dft  ✗ 8/8"],
        ['md_02_homo_lumo_benzene — "HOMO–LUMO gap of benzene at B3LYP/6-31G(d)."',
         "molecular_dft",
         "(no agent)   ·  capability gap"],
    ])
    style_table(tbl, font_size=10)

    # ── Right-bottom: highlight card "Context-aware backend" ──
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.4), Inches(5.95), Inches(5.5), Inches(1.20))
    bg.fill.solid()
    bg.fill.fore_color.rgb = FILL_ALT
    bg.line.color.rgb = ACCENT
    bg.shadow.inherit = False
    add_caption(slide, 7.55, 6.00, 5.2, 0.30,
                "Backend selection is system-context-aware",
                size=12, italic=False, color=ACCENT, align=PP_ALIGN.LEFT)
    add_caption(slide, 7.55, 6.30, 5.2, 0.85,
                "CHGNet picked on battery / cathode prompts (mlip_02).  "
                "16/16 identical picks across both sonnet-4-5 and "
                "opus-4-6 in the variability test.  The agent is reasoning "
                "about training-distribution fit, not defaulting.",
                size=10, italic=True, color=DARK, align=PP_ALIGN.LEFT)

    # ── Bottom-left: by-scale + findings ──
    add_caption(slide, 0.5, 5.95, 6.8, 0.28,
                "Findings",
                size=13, italic=False, color=ACCENT, align=PP_ALIGN.LEFT)
    # ── Variability sweep table ──
    add_caption(slide, 0.5, 5.95, 6.8, 0.28,
                "Variability sweep   (3 prompts × 8 trials × 2 models — "
                "stability = fraction matching modal pick)",
                size=11, italic=False, color=ACCENT, align=PP_ALIGN.LEFT)
    vtbl = slide.shapes.add_table(
        4, 3,
        Inches(0.5), Inches(6.25), Inches(6.8), Inches(1.15)).table
    for c, w in enumerate([3.0, 1.9, 1.9]):
        vtbl.columns[c].width = Inches(w)
    fill_table(vtbl, [
        ["prompt   (difficulty)",         "sonnet-4-5  (8/8)",   "opus-4-6  (8/8)"],
        ["pd_01_lattice_cu   (easy)",     "periodic_dft/vasp ✓", "periodic_dft/vasp ✓"],
        ["mlip_02_battery_cathode  (med)","MLIP/chgnet ✓",       "MLIP/chgnet ✓"],
        ["amb_03_melting_point_cu (hard)","periodic_dft/vasp ✗", "MLIP/mace ✓"],
    ])
    style_table(vtbl, font_size=9)

    add_notes(slide, """
    What this slide shows:
    benchmark/test_router.py — 24 PNNL-relevant prompts fed to
    SimulationRouter, scored on top-1 (scale, engine) accuracy.
    Results under benchmark/outputs/test_router/<model>/.

    Why this matters:
    A wrong scale pick = wrong tool for the problem.  Pick DFT for a
    melting point and you can't sample the liquid; pick a generic
    MLIP for a battery cathode and you miss the dominant chemistry.
    The router is the first decision in every workflow.

    Single-shot headline:
    All three Claude models tested — sonnet-4-5, opus-4-6, opus-4-7 —
    hit 100 % joint accuracy on the 21 scorable queries.  Three
    molecular_dft prompts (gas-phase quantum chemistry) are excluded as
    capability gaps because no MolecularDFTAgent exists in the codebase
    yet.

    Single-shot accuracy isn't the whole story — that's the
    methodological point.

    Variability test:
    Sampled 3 prompts (easy / medium / hard) × 8 trials × 2 models
    (sonnet-4-5 and opus-4-6).
      - pd_01_lattice_cu (easy)          : both 8/8 periodic_dft/vasp
      - mlip_02_battery_cathode (medium) : both 8/8 CHGNet
      - amb_03_melting_point_cu (hard)   : the interesting one

    mlip_02 confirms the qualitative finding on the slide:
    both models picked CHGNet every single trial — 16 identical picks
    across two model architectures.  The context-aware-backend story
    isn't a one-off; it's how these models stably reason about
    training-distribution fit.

    amb_03 is the failure mode the test caught:
    Single-shot, all three models picked MLIP (sonnet→CHGNet, both
    opus→MACE) and passed.  But the 8-trial variability test showed
    sonnet picks periodic_dft/vasp 8 times out of 8.  That sonnet
    single-shot pass was a minority outcome from its distribution.

    Sonnet's stable preference on this prompt is DFT — and DFT alone
    can't actually solve the problem.  You need MD (classical or MLIP)
    to sample the liquid phase for two-phase coexistence or
    thermodynamic integration.  Under a tight rubric, sonnet is stably
    wrong on this prompt.  Opus-4-6 picked MLIP/MACE 8/8 — correct on
    every trial.

    Two takeaways for the audience:

    (1) Sonnet has a real router-quality issue on melting-point-style
        prompts: it picks DFT despite DFT being insufficient for the
        physics.  Opus models don't show this failure mode in the
        cases we sampled.

    (2) Single-shot accuracy is necessary but not sufficient for
        evaluating router quality on ambiguous prompts.  A 100 %
        single-shot score buries a stable failure mode that 8-trial
        sampling exposes.  This is the value-add of the methodology.

    Capability-gap framing:
    The three molecular_dft prompts (dipole of formaldehyde, HOMO-LUMO
    of benzene, CCSD(T) CO dimer binding) get excluded from accuracy
    because there's no MolecularDFTAgent in the codebase.  Different
    models handle them differently — sonnet refuses the most explicit
    cluster-code case (returns None/None); opus models fall back to a
    periodic-DFT-in-a-box workaround.  Sonnet is more honest about
    missing capability; opus tries harder.  Either is defensible;
    worth flagging to PNNL users.

    Honest verdict on model comparison:
    On the bulk of the suite (20 of 21 scorable queries) all three
    Claude models behave the same.  Variability sampling caught one
    stable failure mode in sonnet (DFT on melting-point) that
    opus-4-6 doesn't share.  N=1 — drawing a sweeping
    "opus > sonnet" conclusion would be overclaim.  This run
    justifies "watch sonnet on open-ended thermodynamics prompts".
    """)


def slide_gen_and_run(prs: Presentation) -> None:
    """Slide 6/6: post-router stages — INCAR generation, DFT cells, MLIP→MD."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Generation + run  —  INCAR quality, DFT cells, MLIP→MD")
    add_subtitle(slide,
                 "Once the router picks the tool, can the agent produce a "
                 "correct input deck, run it, and read the result?  Three "
                 "tests at three stages of the pipeline.")

    # ── Top: 3 headline number cards ──
    def _stat_card(x: float, big: str, label: str, accent=ACCENT) -> None:
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(1.7), Inches(4.0), Inches(1.0))
        bg.fill.solid()
        bg.fill.fore_color.rgb = ORCH_FILL
        bg.line.color.rgb = accent
        bg.shadow.inherit = False
        box = slide.shapes.add_textbox(
            Inches(x), Inches(1.78), Inches(4.0), Inches(0.55))
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = big
        run = p.runs[0]
        run.font.size = Pt(32); run.font.bold = True; run.font.color.rgb = accent
        box2 = slide.shapes.add_textbox(
            Inches(x), Inches(2.30), Inches(4.0), Inches(0.4))
        p2 = box2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        p2.text = label
        run2 = p2.runs[0]
        run2.font.size = Pt(11); run2.font.italic = True
        run2.font.color.rgb = DARK

    _stat_card(0.5, "6 / 9", "DFT cells pass within ±1.3 % of experiment")
    _stat_card(4.7, "5 / 5", "MLIP→MD trajectories complete  (MACE-MP-0)",
               accent=ACCENT)
    _stat_card(8.9, "46 % vs 21 %",
               "INCAR typo rate   sonnet vs opus-4-6",
               accent=SKILL_LINE)

    # ── Middle: "Why the typo rate matters" framing box ──
    add_caption(slide, 0.5, 2.95, 12.3, 0.28,
                "Why the typo rate matters",
                size=12, italic=False, color=ACCENT, align=PP_ALIGN.LEFT)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(3.25), Inches(12.3), Inches(0.90))
    bg.fill.solid()
    bg.fill.fore_color.rgb = FILL_ALT
    bg.line.color.rgb = SKILL_LINE
    bg.shadow.inherit = False
    add_caption(slide, 0.65, 3.32, 12.0, 0.80,
                "VASP silently ignores unknown INCAR tags — a one-letter typo "
                "like ISPN-for-ISPIN disables spin polarization with no error.  "
                "Run completes; physics is wrong.  Fe BCC came back at 2.756 Å "
                "(NM-Fe lattice) the first pass.  We added a pre-submit syntax "
                "check (pymatgen) that catches every high-confidence typo "
                "before sbatch.  Fe rerun: 2.831 Å, magmom 2.20 μB — passes.",
                size=11, italic=False, color=DARK, align=PP_ALIGN.LEFT)

    # ── DFT cells — one-line summary above the two tables ──
    add_caption(slide, 0.5, 4.35, 12.3, 0.28,
                "DFT cells  (Deception, 9 systems):   Cu · Si · C · MgO · LiCoO₂ · UO₂ "
                "all ≤ 1.3 % of experiment.   Fe passes after pre-submit validator.   "
                "Pt(111)+CO and TiO₂(101)+H₂O slabs re-running with extended wall-time.",
                size=10, italic=False, color=DARK, align=PP_ALIGN.LEFT)

    # ── Bottom-left: INCAR variability table ──
    add_caption(slide, 0.5, 4.75, 6.3, 0.25,
                "INCAR variability   (3 prompts × 8 trials × 2 models, generate-only)",
                size=11, italic=False, color=ACCENT, align=PP_ALIGN.LEFT)
    itbl = slide.shapes.add_table(
        5, 5,
        Inches(0.5), Inches(5.05), Inches(6.3), Inches(2.10)).table
    for c, w in enumerate([2.2, 1.05, 1.05, 1.05, 1.05]):
        itbl.columns[c].width = Inches(w)
    fill_table(itbl, [
        ["prompt",            "sonnet typo", "sonnet stab", "opus typo", "opus stab"],
        ["fe_bcc_magnetic",   "0.50  (ISPN×4)", "1.00",      "0.00",     "1.00"],
        ["uo2_dftU",          "0.50  (LMMAXMIX×2 …)", "1.00", "0.62  (LMIXFIXED×3 …)", "1.00"],
        ["pt111_co_dipole",   "0.38  (ISPN×3)", "0.62",      "0.00",     "0.88"],
        ["mean",              "0.46",        "0.88",        "0.21",     "0.96"],
    ])
    style_table(itbl, font_size=9)

    # ── Bottom-right: MLIP systems table ──
    add_caption(slide, 7.0, 4.75, 5.8, 0.25,
                "MLIP → MD pipeline test   (Deception, ASE runner, 1000-step NVT @ 300 K)",
                size=11, italic=False, color=ACCENT, align=PP_ALIGN.LEFT)
    mtbl = slide.shapes.add_table(
        7, 3,
        Inches(7.0), Inches(5.05), Inches(5.8), Inches(2.10)).table
    for c, w in enumerate([2.6, 2.2, 1.0]):
        mtbl.columns[c].width = Inches(w)
    fill_table(mtbl, [
        ["system",            "chemistry class",        "wall time"],
        ["Fe BCC",            "ferromagnetic metal",    "35.9 s"],
        ["Cu FCC",            "noble metal",            "35.5 s"],
        ["Si diamond",        "covalent semiconductor", "41.0 s"],
        ["MgO rocksalt",      "ionic insulator",        "36.4 s"],
        ["LiCoO₂ layered",    "battery cathode",        "40.0 s"],
        ["water_box",         "polar liquid (skipped — Packmol)", "—"],
    ])
    style_table(mtbl, font_size=9)

    add_notes(slide, """
    What this slide shows:
    Three tests covering everything downstream of the router —
    INCAR generation (test_incar_variability), DFT cell relaxation
    (test_dft), and MLIP→MD handoff (test_mlip).  All three live in
    benchmark/, all three run on Deception, results in
    benchmark/outputs/<test>/.

    Why this matters:
    Picking the right scale (router slide) is necessary but not
    sufficient.  The agent then has to produce a correct input deck
    for that engine, run it, and parse the result.  Each of those
    stages is its own failure surface.

    DFT cells (benchmark/test_dft, 9 systems):
    Pass-rate 6/9, mean score 0.74 (±2.5 % lattice-constant band).
    Pass:  Cu, Si, C-diamond, MgO, LiCoO₂, UO₂ — all within 1.3 %
    of experiment.  Fe BCC failed the first pass — see Fe story below.
    Pt(111)+CO and TiO₂(101)+H₂O adsorption fragments did not finish
    in the wall-time window; re-running with extended wall-time.

    The MgO pass is itself a small validation story:  previously
    failed silently because ase.io.write(vasp) emitted interleaved
    species without sort=True, so the POSCAR/POTCAR ordering
    disagreed.  Fix landed in 655aa7e.

    The Fe BCC story — pre-submit validator earned its keep:
    First Fe run came back at 2.756 Å vs experimental 2.866 Å — a
    −3.9 % miss.  PBE typically *over*-estimates Fe (~2.83 Å for FM
    Fe), so the magnitude and direction were both suspicious.  The
    INCAR the agent emitted:  ``ISPN = 2`` instead of ``ISPIN = 2``.
    VASP silently ignores unknown tags — the run completed cleanly
    with no error, but spin polarisation was off, so we got the
    *non*-magnetic Fe lattice constant (~2.76 Å, exactly what landed).
    OUTCAR confirms magnetic moment ≈ 0.

    This is the canonical "VASP-converged but physics-wrong" failure
    mode.  Today's commit (6d7206c on validation-suite) plumbs a
    pre-submit syntax check into PeriodicDFTAgent:  every generated
    INCAR goes through pymatgen's Incar.check_params() before
    submission.  High-confidence typos auto-renamed via difflib match
    against the canonical VASP tag list.  Same engine-neutral shape
    as DeployedPotential — pymatgen lives behind a per-engine module
    (vasp_input_validator.py), no cross-engine code knows about
    INCAR tags.  LAMMPS and GROMACS validators will follow the same
    shape once they land.

    Fe re-ran through this path and now passes: a = 2.831 Å
    (Δ −1.24 % vs experiment, within ±2.5 % band), magnetic moment
    2.20 μB / atom (vs expected 2.22).  Validator earned its keep,
    end-to-end.

    INCAR variability test (benchmark/test_incar_variability):
    Same N-shot methodology as the router variability test, applied
    one stage downstream.  3 prompts × 8 trials × 2 models, all
    generate-only (no cluster):
      * fe_bcc_magnetic    — ISPIN + MAGMOM choices
      * uo2_dftU           — LDAUL/LDAUU/LDAUJ + AFM init
      * pt111_co_dipole    — IDIPOL/LDIPOL + asymmetric slab
    Each trial scored on (a) typo rate via pymatgen Incar.check_params,
    (b) physics-tuple stability across (ISPIN, ISIF, ISMEAR, ENCUT-band,
    NSW>0, MAGMOM-set, LDAU-set, IDIPOL/LDIPOL-set).

    sonnet-4-5  (8 × 3 = 24 trials):
      - mean typo rate 0.46     -- ISPN-for-ISPIN in 8/24 trials (33 %),
                                   always auto-fixed at high confidence
      - mean physics stability 0.88
      - UO₂: 2 confabulated tags (LDAALPHA, LDMIXINGLINEAR) flagged
        low-confidence and NOT auto-renamed — correct conservative
        behavior, those have no close VASP-tag match

    opus-4-6  (8 × 3 = 24 trials):
      - mean typo rate 0.21     -- Fe 0/8, Pt 0/8.  UO₂ 5/8 with a
                                   *different* confabulated family
                                   (LMIXFIXED, LMIXFMT, LMIXMAG)
      - mean physics stability 0.96

    Read:  opus is the cleaner generator overall (Fe and Pt clean
    across 16 trials); both models stumble on UO₂ DFT+U where
    "mixing-related" tags get confabulated.  All sonnet typos are
    auto-fixed at high confidence; opus's UO₂ typos are correctly
    flagged but not silently renamed (no close-enough VASP match).
    Pre-submit validator is the right safety net for both models,
    not just for the typo-prone one.

    MLIP→MD (benchmark/test_mlip, 5 cells + 1 skipped):
    All five solid-state cells deployed MACE-MP-0 and completed a
    1000-step NVT @ 300 K MD trajectory on a single A100, ~36–41 s
    per cell.  This is the first end-to-end MLIP→MD demonstration on
    the stack:  MLIPAgent emits a DeployedPotential descriptor and
    MDSimulationAgent (ASE runner) consumes it without importing any
    MLIP code itself — the engine-neutral N+M contract from
    CLAUDE.md.  All five picked MACE because mace-mp-0 covers every
    element in those cells; the SimulationRouter's CHGNet preference
    for LiCoO₂ lives on a different code path and doesn't constrain
    the MLIPAgent's own backend choice.  water_box was skipped — it's
    a Packmol spec and test_mlip doesn't yet materialise liquid boxes
    from molecular fragments (known gap).
    The runs used the ASE runner rather than LAMMPS+MACE; the LAMMPS
    build is a separate work-stream and was deferred to keep the
    panel moving.
    """)


def main() -> None:
    prs = Presentation()
    set_widescreen(prs)

    slide_title(prs)
    slide_stack(prs)
    slide_pipelines(prs)
    slide_autonomy(prs)
    slide_router(prs)
    slide_gen_and_run(prs)

    out = "scilink_sim_agents_overview.pptx"
    prs.save(out)
    print(f"wrote {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
