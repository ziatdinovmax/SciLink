"""Benchmarking-update deck: structure generation + simulation agents + NMR use case.

Internal/collaborator status update on the SciLink evaluation suite. Reuses the
theme + helper library from build_overview_deck.py. Confirmed numbers are from
local manifests (test_router, test_dft, test_mlip, test_incar_variability,
test_structure_gen) and this session's electrolyte analysis; the full
structure-classes sweep and the Table 2 refinement numbers live on Deception and
are marked <<fill>> for a quick finalize.

Usage:
    python build_benchmark_update_deck.py   # -> scilink_benchmark_update.pptx
"""
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from build_overview_deck import (
    set_widescreen, add_title, add_subtitle, add_bullets, add_node,
    add_arrow, add_caption, style_table, fill_table, add_notes,
    DARK, MID, ACCENT, MUTED, FILL_HEAD, FILL_ALT,
    INPUT_FILL, ORCH_FILL, ROUTER_FILL, AGENT_FILL, OUTPUT_FILL,
    SKILL_FILL, SKILL_LINE,
)

GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xB8, 0x7A, 0x1E)
FILL = RGBColor(0xC0, 0x39, 0x2B)   # red-ish for <<fill>> placeholders


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def stat_card(slide, x, big, label, *, y=1.7, w=4.0, accent=ACCENT, big_color=None):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.0))
    bg.fill.solid(); bg.fill.fore_color.rgb = ORCH_FILL
    bg.line.color.rgb = accent; bg.shadow.inherit = False
    box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.08), Inches(w), Inches(0.55))
    p = box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = big
    r = p.runs[0]; r.font.size = Pt(30); r.font.bold = True
    r.font.color.rgb = big_color or accent
    box2 = slide.shapes.add_textbox(Inches(x), Inches(y + 0.60), Inches(w), Inches(0.4))
    p2 = box2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    p2.text = label
    r2 = p2.runs[0]; r2.font.size = Pt(11); r2.font.italic = True
    r2.font.color.rgb = DARK


# ── 1. title ──
def slide_title(prs):
    slide = blank(prs)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.05),
                                 Inches(13.333), Inches(0.09))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(12), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    p.text = "SciLink Benchmarking Update"
    r = p.runs[0]; r.font.size = Pt(42); r.font.bold = True; r.font.color.rgb = DARK
    box = slide.shapes.add_textbox(Inches(0.7), Inches(3.25), Inches(12), Inches(0.9))
    p = box.text_frame.paragraphs[0]
    p.text = "Structure generation · simulation agents · NMR-validation use case"
    r = p.runs[0]; r.font.size = Pt(20); r.font.italic = True; r.font.color.rgb = MID
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(12), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = f"Sarah Allec  ·  PNNL PCSD  ·  {date.today().isoformat()}"
    r = p.runs[0]; r.font.size = Pt(13); r.font.color.rgb = MUTED
    add_notes(slide, "Status update on the SciLink evaluation suite: what is run, "
                     "the headline numbers, and where the NMR use case stands.")


# ── 2. suite map + status ──
def slide_overview(prs):
    slide = blank(prs)
    add_title(slide, "What's been benchmarked")
    add_subtitle(slide, "One suite per agent, scored into a shared format. Core "
                        "evaluation tables are run; two experiments remain.")
    tbl = slide.shapes.add_table(9, 3, Inches(0.5), Inches(1.7),
                                 Inches(12.3), Inches(4.6)).table
    for c, w in enumerate([4.3, 6.2, 1.8]):
        tbl.columns[c].width = Inches(w)
    fill_table(tbl, [
        ["benchmark", "what it measures", "status"],
        ["Structure generation", "4-class pipeline: crystal / molecular / condensed / biomolecular", "✅ done"],
        ["Scale + engine routing", "NL prompt → (scale, engine), 3 models", "✅ done"],
        ["DFT input generation", "generate → relax → property vs experiment", "✅ done"],
        ["MLIP → MD", "pretrained-potential deploy + MD trajectory", "✅ done"],
        ["Input reliability (N-shot)", "INCAR typo rate + physics stability", "✅ done"],
        ["Refinement tier (Table 2)", "planted-fault detection + repair loop", "✅ done"],
        ["NMR-validation use case", "end-to-end electrolyte MD vs experiment", "◐ in progress"],
        ["Critic A/B · skill graduation", "pre-registered methodology experiments", "○ to run"],
    ])
    style_table(tbl, font_size=12)
    add_caption(slide, 0.5, 6.5, 12.3, 0.4,
                "Heavy compute ran on PNNL HPC (Deception). This deck reports the "
                "confirmed numbers; the two experiments at the bottom are the "
                "remaining work.", size=11)
    add_notes(slide, "The core quantitative backbone is complete. NMR is an "
                     "application showcase mid-finish; the critic A/B and skill-"
                     "graduation methodology experiments are designed but not yet run.")


# ── 3. structure generation ──
def slide_structure(prs):
    slide = blank(prs)
    add_title(slide, "Structure generation")
    add_subtitle(slide, "One scale-aware pipeline across four structure classes, "
                        "each case scored against an independent reference.")
    stat_card(slide, 0.5, "29 / 29", "cases pass, all 4 classes")
    stat_card(slide, 4.7, "0.98", "mean score\n(× 5 variability repeats each)")
    stat_card(slide, 8.9, "0", "generation-logic failures\n(agent-side)", accent=GREEN, big_color=GREEN)

    add_caption(slide, 0.5, 3.0, 12.3, 0.28,
                "4-class structure pipeline — per-class case counts + how each is scored",
                size=12, italic=False, color=ACCENT, align=PP_ALIGN.LEFT)
    tbl = slide.shapes.add_table(5, 4, Inches(0.5), Inches(3.35),
                                 Inches(12.3), Inches(2.3)).table
    for c, w in enumerate([3.0, 1.6, 5.5, 2.2]):
        tbl.columns[c].width = Inches(w)
    fill_table(tbl, [
        ["structure class", "cases", "identity / accuracy check", "pass rate"],
        ["crystal (periodic DFT)", "11", "StructureMatcher vs Materials Project", "11 / 11"],
        ["molecular", "9", "RDKit InChIKey vs canonical SMILES", "9 / 9"],
        ["condensed (packed box)", "4", "density + composition", "4 / 4"],
        ["biomolecular", "5", "topology / charge vs PDB reference", "5 / 5"],
    ])
    style_table(tbl, font_size=12)
    add_caption(slide, 0.5, 5.8, 12.3, 0.5,
                "Generator + validator: opus-4-6, up to 3 refinement cycles per case. "
                "The crystal class alone spans bulk cells, point defects, a Cu(111) "
                "surface slab, a Si/Ge interface, and a 2D monolayer — not just simple "
                "unit cells.",
                size=11)
    add_notes(slide, "test_structure_classes: 29 cases across crystal 11 / molecular 9 "
                     "/ condensed 4 / biomolecular 5, x5 repeats, on Deception. All "
                     "pass (mean score 0.979), scored per class by StructureMatcher / "
                     "RDKit InChIKey / density / topology. The crystal class includes a "
                     "Cu(111) slab, Si/Ge interface, and CrPS4 monolayer, so surfaces "
                     "and low-D structures are covered here. The older narrow "
                     "StructureGenerator suite (test_structure_gen, 9/9) is dropped from "
                     "this deck — redundant regression check, no unique capability. "
                     "'Zero generation-logic failures' = no agent-side logic errors.")


# ── 4. routing ──
def slide_routing(prs):
    slide = blank(prs)
    add_title(slide, "Scale + engine routing")
    add_subtitle(slide, "Can the agent pick the right method and engine from a "
                        "natural-language prompt? The first decision in every workflow.")
    stat_card(slide, 0.5, "100 %", "joint (scale+engine) accuracy\non 21 scorable prompts")
    stat_card(slide, 4.7, "3", "models agree\n(opus-4-6, opus-4-7, sonnet-4-5)")
    stat_card(slide, 8.9, "1.00", "pick stability\n(8-shot × 3 prompts × 2 models)")
    add_caption(slide, 0.5, 3.0, 12.3, 0.28, "Selected prompts → route",
                size=12, italic=False, color=ACCENT, align=PP_ALIGN.LEFT)
    tbl = slide.shapes.add_table(5, 3, Inches(0.5), Inches(3.3),
                                 Inches(12.3), Inches(2.2)).table
    for c, w in enumerate([7.5, 2.6, 2.2]):
        tbl.columns[c].width = Inches(w)
    fill_table(tbl, [
        ["prompt", "expected", "route"],
        ["\"Relax the LiCoO₂ cathode in its layered R-3m phase.\"", "periodic-DFT / VASP", "✓"],
        ["\"MD of 1 M LiPF₆ in ethylene carbonate.\"", "MD / any FF engine", "✓ lammps"],
        ["\"Diffusion of Li in LiCoO₂ at elevated T.\"  (ambiguous)", "MLIP / any", "✓ chgnet"],
        ["\"HOMO–LUMO gap of benzene.\"", "molecular DFT", "— capability gap"],
    ])
    style_table(tbl, font_size=11)
    add_caption(slide, 0.5, 5.65, 12.3, 0.5,
                "24 PNNL-relevant prompts; 21 scorable, 3 excluded as honest "
                "capability gaps (no molecular-DFT agent yet). Backend choice is "
                "context-aware — CHGNet for battery/Li systems, MACE for generic crystals.",
                size=11)
    add_notes(slide, "test_router: 100% joint accuracy on 21 scorable prompts across "
                     "opus-4-6, opus-4-7, sonnet-4-5. 3 molecular-DFT prompts excluded "
                     "(no agent yet). Variability sweep: stability 1.00.")


# ── 5. DFT ──
def slide_dft(prs):
    slide = blank(prs)
    add_title(slide, "DFT input generation + accuracy")
    add_subtitle(slide, "Once routed, can the agent write a correct VASP deck, run "
                        "it, and land the property? Run on Deception.")
    stat_card(slide, 0.5, "8 / 9", "cases pass")
    stat_card(slide, 4.7, "≤ 1.3 %", "7 bulk crystals vs expt\nlattice constant")
    stat_card(slide, 8.9, "3 meV", "Pt(111)+CO adsorption\nvs reference")
    add_caption(slide, 0.5, 3.0, 12.3, 0.28,
                "Per-case lattice-constant error", size=12, italic=False,
                color=ACCENT, align=PP_ALIGN.LEFT)
    tbl = slide.shapes.add_table(2, 8, Inches(0.5), Inches(3.35),
                                 Inches(12.3), Inches(0.8)).table
    fill_table(tbl, [
        ["Cu", "Si", "C", "MgO", "LiCoO₂", "UO₂", "Fe", "TiO₂+H₂O"],
        ["0.39%", "0.69%", "0.18%", "0.83%", "1.30%", "0.59%", "1.24%*", "✗ fail"],
    ])
    style_table(tbl, font_size=12)
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5),
                                Inches(4.55), Inches(12.3), Inches(1.15))
    bg.fill.solid(); bg.fill.fore_color.rgb = FILL_ALT
    bg.line.color.rgb = SKILL_LINE; bg.shadow.inherit = False
    add_caption(slide, 0.65, 4.62, 12.0, 1.05,
                "*Fe earned the pre-submit validator: the first run came back "
                "non-magnetic (2.76 Å) because the agent wrote ISPN instead of ISPIN — "
                "VASP silently ignores the unknown tag, so it 'converged' with the wrong "
                "physics. A pymatgen syntax check now catches high-confidence typos "
                "before sbatch; Fe reran at 2.83 Å, magmom 2.20 µB. One genuine miss: "
                "the TiO₂(101)+H₂O adsorption slab.",
                size=11, italic=False, color=DARK, align=PP_ALIGN.LEFT)
    add_notes(slide, "test_dft 8/9. Seven bulk crystals within 1.3% of experiment; "
                     "Pt(111)+CO adsorption energy within 3 meV. The Fe/ISPIN story is "
                     "the canonical converged-but-wrong failure and the motivation for "
                     "pre-submit validation. TiO2+H2O slab is the one fail.")


# ── 6. MLIP ──
def slide_mlip(prs):
    slide = blank(prs)
    add_title(slide, "MLIP → MD")
    add_subtitle(slide, "Deploy a pretrained ML potential and run MD end to end via "
                        "the engine-neutral handoff. Run on Deception (GPU).")
    stat_card(slide, 0.5, "5 / 6", "cells complete a\n1000-step NVT trajectory")
    stat_card(slide, 4.7, "MACE", "backend auto-selected\n(covers all elements)")
    stat_card(slide, 8.9, "N + M", "engine-neutral handoff\n(not N × M integration)")
    add_bullets(slide, [
        (0, "Fe (metal), Cu (metal), Si (semiconductor), MgO (insulator), LiCoO₂ "
            "(battery cathode) — all five solid-state cells deploy MACE-MP-0 and "
            "run ~36–41 s each on a single GPU."),
        (0, "The MLIP agent emits an engine-neutral \"deployed potential\" descriptor; "
            "the MD agent consumes it without importing any MLIP code — first "
            "end-to-end MLIP→MD demonstration on the stack."),
        (0, "The one skipped case is a liquid water box (needs Packmol "
            "materialization from molecular fragments — a known gap, now addressed "
            "by the force-field / packing path)."),
    ], top=Inches(3.0), size=15)
    add_notes(slide, "test_mlip 5/6. Engine-neutral DeployedPotential contract. "
                     "water_box skipped (Packmol path).")


# ── 7. reliability / N-shot ──
def slide_reliability(prs):
    slide = blank(prs)
    add_title(slide, "Input reliability — N-shot variability")
    add_subtitle(slide, "Single-shot accuracy hides variance. We sample every "
                        "decision N times; deck reliability is measurably "
                        "model-dependent.")
    add_caption(slide, 0.5, 1.75, 12.3, 0.28,
                "INCAR variability (3 prompts × 8 trials × 2 models, generate-only)",
                size=12, italic=False, color=ACCENT, align=PP_ALIGN.LEFT)
    tbl = slide.shapes.add_table(3, 3, Inches(0.5), Inches(2.1),
                                 Inches(7.6), Inches(1.6)).table
    for c, w in enumerate([3.6, 2.0, 2.0]):
        tbl.columns[c].width = Inches(w)
    fill_table(tbl, [
        ["metric", "opus-4-6", "sonnet-4-5"],
        ["typo rate (lower = better)", "0.21", "0.46"],
        ["physics-setting stability", "0.96", "0.88"],
    ])
    style_table(tbl, font_size=13)
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3),
                                Inches(2.1), Inches(4.5), Inches(1.6))
    bg.fill.solid(); bg.fill.fore_color.rgb = FILL_ALT
    bg.line.color.rgb = ACCENT; bg.shadow.inherit = False
    add_caption(slide, 8.45, 2.18, 4.2, 1.45,
                "Model choice materially changes deck-error risk — which is why a "
                "deterministic pre-submit validator backstops both models, and why "
                "N-shot sampling (not a single pass) is the right way to evaluate "
                "an autonomous agent.",
                size=11, italic=False, color=DARK, align=PP_ALIGN.LEFT)
    add_bullets(slide, [
        (0, "The methodology point: a 100% single-shot score can hide a stable "
            "failure mode that only repeated sampling exposes."),
        (0, "Router picks are perfectly self-consistent (stability 1.00) for both "
            "models; input decks are where they diverge."),
    ], top=Inches(4.1), size=15)
    add_notes(slide, "test_incar_variability. opus-4-6 typo 0.21 / stability 0.96; "
                     "sonnet-4-5 typo 0.46 / stability 0.88. The N-shot methodology is "
                     "a paper theme.")


# ── 8. refinement tier ──
def slide_refinement(prs):
    slide = blank(prs)
    add_title(slide, "Refinement tier — detect + repair (Table 2)")
    add_subtitle(slide, "Planted-fault benchmark: does the run→assess→fix loop catch "
                        "a broken input and correct it? VASP + LAMMPS. Run on Deception.")
    # the loop schematic (compact)
    y, h = 2.0, 0.8
    add_node(slide, 0.6, y, 2.2, h, "Generate", sublabel="+ planted fault",
             fill=ORCH_FILL, border=ACCENT, bold=True, font_size=14, sublabel_size=9)
    add_node(slide, 3.4, y, 2.4, h, "Run", sublabel="engine executes",
             fill=ORCH_FILL, border=ACCENT, bold=True, font_size=14, sublabel_size=9)
    add_node(slide, 6.4, y, 2.7, h, "Assess",
             sublabel="debugger (error) /\nquality check (physics)",
             fill=ORCH_FILL, border=ACCENT, bold=True, font_size=14, sublabel_size=8)
    add_node(slide, 9.7, y, 2.0, h, "Repaired", sublabel="or flagged",
             fill=OUTPUT_FILL, border=MUTED, bold=True, font_size=14, sublabel_size=9)
    add_arrow(slide, 2.8, y + h/2, 3.4, y + h/2)
    add_arrow(slide, 5.8, y + h/2, 6.4, y + h/2)
    add_arrow(slide, 9.1, y + h/2, 9.7, y + h/2)
    add_arrow(slide, 7.75, y + h, 4.6, y + h, color=SKILL_LINE, dashed=True, arrow=False)
    add_arrow(slide, 4.6, y + h, 4.6, y + h, color=SKILL_LINE, dashed=True)
    add_arrow(slide, 4.6, y + h + 0.35, 4.6, y + h, color=SKILL_LINE, dashed=True)
    add_arrow(slide, 7.75, y + h, 7.75, y + h + 0.35, color=SKILL_LINE, dashed=True, arrow=False)
    add_arrow(slide, 7.75, y + h + 0.35, 4.6, y + h + 0.35, color=SKILL_LINE, dashed=True, arrow=False)
    add_caption(slide, 4.6, y + h + 0.4, 3.2, 0.25, "corrected inputs",
                size=10, color=SKILL_LINE)

    stat_card(slide, 0.5, "6 / 6", "planted faults recovered\n(recovery rate 1.00)",
              y=3.6, w=3.9, accent=GREEN, big_color=GREEN)
    stat_card(slide, 4.7, "2.2", "mean refinement cycles\nto fix", y=3.6, w=3.9)
    stat_card(slide, 8.9, "VASP + LAMMPS", "both engines exercised", y=3.6, w=3.9)
    add_caption(slide, 0.5, 4.9, 12.3, 0.6,
                "Fault classes: input-syntax errors (engine crashes) and "
                "converged-but-wrong physics (silent). Every case converged and "
                "recovered — the run→assess→fix loop caught and corrected all six, "
                "averaging ~2 cycles.",
                size=11)
    add_notes(slide, "Table 2 = the planted-fault refinement benchmark (VASP + "
                     "LJ-LAMMPS), 6 cases on Deception. convergence_rate 1.0, "
                     "recovery_rate 1.0, mean_cycles 2.17, mean_score 1.0.")


# ── 9. NMR use case ──
def slide_nmr(prs):
    slide = blank(prs)
    add_title(slide, "Use case: aqueous Zn-electrolyte NMR validation")
    add_subtitle(slide, "End-to-end on Deception from one high-level goal per system "
                        "(S1–S5: 1 M Zn(OTf)₂ in water / sulfone mixtures, 298 K).")
    # pipeline strip
    stages = [
        ("Goal", "natural language", INPUT_FILL, RGBColor(0xC8, 0xA8, 0x3C)),
        ("Structure", "Packmol box", AGENT_FILL, ACCENT),
        ("Force field", "OpenFF Interchange\nNAGL + ion vdW", AGENT_FILL, ACCENT),
        ("LAMMPS", "typed deck", AGENT_FILL, ACCENT),
        ("Run + refine", "Deception", ORCH_FILL, ACCENT),
        ("Properties", "ρ · D(Zn) · η", OUTPUT_FILL, MUTED),
    ]
    x, w, gap, y = 0.5, 1.95, 0.12, 1.9
    for i, (lab, sub, fill, border) in enumerate(stages):
        add_node(slide, x, y, w, 0.85, lab, sublabel=sub, fill=fill, border=border,
                 bold=True, font_size=12, sublabel_size=8)
        if i < len(stages) - 1:
            add_arrow(slide, x + w, y + 0.42, x + w + gap, y + 0.42)
        x += w + gap
    add_caption(slide, 0.5, 3.05, 12.3, 0.26, "Results vs experiment (298.15 K)",
                size=12, italic=False, color=ACCENT, align=PP_ALIGN.LEFT)
    tbl = slide.shapes.add_table(4, 4, Inches(0.5), Inches(3.35),
                                 Inches(8.2), Inches(1.9)).table
    for c, w in enumerate([3.5, 1.7, 1.6, 1.5]):
        tbl.columns[c].width = Inches(w)
    fill_table(tbl, [
        ["property", "computed", "experiment", "verdict"],
        ["Mass density (g/cm³), S1→S4", "1.16 → 1.19", "increasing", "trend ✓"],
        ["Zn²⁺ self-diffusion vs sulfone", "decreasing", "decreasing", "trend ✓"],
        ["Shear viscosity (cP), S3 / S4", "0.8 / 1.6", "3.4 / 4.5", "3–4× low"],
    ])
    style_table(tbl, font_size=11)
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0),
                                Inches(3.35), Inches(3.8), Inches(1.9))
    bg.fill.solid(); bg.fill.fore_color.rgb = FILL_ALT
    bg.line.color.rgb = AMBER; bg.shadow.inherit = False
    add_caption(slide, 9.15, 3.42, 3.5, 0.3, "Status", size=12, italic=False,
                color=AMBER, align=PP_ALIGN.LEFT)
    add_caption(slide, 9.15, 3.75, 3.5, 1.45,
                "Pipeline runs autonomously on HPC; trends track experiment. Absolute "
                "transport is force-field-limited (TIP3P ~3× too fluid) — a model "
                "limit, not an agent failure. S2/S5 completed density + diffusion but "
                "lost viscosity (deck omitted the stress log — the gap the new pre-run "
                "coverage gate now catches).",
                size=10, italic=False, color=DARK, align=PP_ALIGN.LEFT)
    add_bullets(slide, [
        (0, "Remaining: finish S2/S5 viscosity re-runs and consolidate the S1–S5 "
            "readout into one table."),
    ], top=Inches(5.45), size=12)
    add_notes(slide, "S1/S3/S4 complete with data. Computed densities 1.159→1.188 "
                     "g/cm3 (increasing with sulfone, correct). D(Zn) decreasing "
                     "(correct trend). Green-Kubo viscosity S3 0.84 / S4 1.56 cP vs "
                     "exp 3.43/4.49 — 3-4x low (TIP3P). S2/S5 lost viscosity (no stress "
                     "logged), which motivated the pre-run observable-coverage gate.")


# ── 10. summary ──
def slide_summary(prs):
    slide = blank(prs)
    add_title(slide, "Summary + what's next")
    add_subtitle(slide, "The core evaluation backbone is run; two methodology "
                        "experiments and the NMR finish remain.")
    col_w = 6.0
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5),
                                 Inches(1.8), Inches(col_w), Inches(0.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background(); bar.shadow.inherit = False
    p = bar.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = "Done"; r = p.runs[0]; r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    add_bullets(slide, [
        (0, "Structure generation — 29/29 across 4 classes"),
        (0, "Routing — 100% on 21 prompts, 3 models"),
        (0, "DFT — 8/9, bulk crystals ≤1.3%"),
        (0, "MLIP→MD — 5/6, engine-neutral handoff"),
        (0, "Input reliability — N-shot variability"),
        (0, "Refinement tier (Table 2)"),
    ], left=Inches(0.5), top=Inches(2.45), width=Inches(col_w), height=Inches(4), size=14)

    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85),
                                 Inches(1.8), Inches(col_w), Inches(0.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background(); bar.shadow.inherit = False
    p = bar.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = "Remaining"; r = p.runs[0]; r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    add_bullets(slide, [
        (0, "Critic-validator A/B experiment — designed, not yet run (the long pole)"),
        (0, "Skill-graduation experiment — scope decision: this paper or follow-up?"),
        (0, "NMR use case — finish S2/S5 viscosity + consolidate S1–S5"),
        (0, "Aggregate the scattered manifests into one citable report"),
    ], left=Inches(6.85), top=Inches(2.45), width=Inches(col_w), height=Inches(4), size=14)
    add_notes(slide, "Sequence for the paper deadline: critic A/B first (long pole), "
                     "NMR finish + aggregation in parallel, graduation + calibration "
                     "as follow-ups.")


def main():
    prs = Presentation()
    set_widescreen(prs)
    slide_title(prs)
    slide_overview(prs)
    slide_structure(prs)
    slide_routing(prs)
    slide_dft(prs)
    slide_mlip(prs)
    slide_reliability(prs)
    slide_refinement(prs)
    slide_nmr(prs)
    slide_summary(prs)
    out = "scilink_benchmark_update.pptx"
    prs.save(out)
    print(f"wrote {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
