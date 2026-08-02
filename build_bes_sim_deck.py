"""Generate a standalone ~7-slide section on SciLink's simulation side for a
BES Autonomy workshop (external audience).

Focus: engine-neutral atomistic simulation agents that set up, run, validate,
and refine DFT / classical-MD / MLIP jobs on PNNL HPC (Deception, SLURM), with
recent work foregrounded — the OpenFF force-field path, pre-run validation
gates, and an end-to-end aqueous-electrolyte campaign.

Reuses the theme + helper library from build_overview_deck.py so styling stays
on-brand. Numbers are grounded in benchmark/outputs/ and the electrolyte runs.

Usage:
    python build_bes_sim_deck.py
    # -> scilink_bes_sim_deck.pptx
"""
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from build_overview_deck import (
    set_widescreen, add_title, add_subtitle, add_bullets, add_node,
    add_arrow, add_caption, style_table, fill_table, add_notes,
    DARK, MID, ACCENT, MUTED, FILL_HEAD, FILL_ALT,
    INPUT_FILL, ORCH_FILL, ROUTER_FILL, AGENT_FILL, OUTPUT_FILL,
    SKILL_FILL, SKILL_LINE,
)

GREEN = RGBColor(0x2E, 0x7D, 0x32)   # success
AMBER = RGBColor(0xB8, 0x7A, 0x1E)   # challenge / caution


# ── shared card helper ────────────────────────────────────────────
def stat_card(slide, x, big, label, *, y=1.7, w=4.0, accent=ACCENT):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.0))
    bg.fill.solid()
    bg.fill.fore_color.rgb = ORCH_FILL
    bg.line.color.rgb = accent
    bg.shadow.inherit = False
    box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.08), Inches(w), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = big
    r = p.runs[0]
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = accent
    box2 = slide.shapes.add_textbox(Inches(x), Inches(y + 0.60), Inches(w), Inches(0.4))
    p2 = box2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.text = label
    r2 = p2.runs[0]
    r2.font.size = Pt(11); r2.font.italic = True; r2.font.color.rgb = DARK


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── 1. section title ──────────────────────────────────────────────
def slide_title(prs):
    slide = blank(prs)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.05), Inches(13.333), Inches(0.09))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.75), Inches(12), Inches(1.3))
    p = box.text_frame.paragraphs[0]
    p.text = "SciLink: Autonomous Atomistic Simulation on HPC"
    r = p.runs[0]; r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = DARK

    box = slide.shapes.add_textbox(Inches(0.7), Inches(3.25), Inches(12), Inches(1.0))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ("Engine-neutral DFT / MD / MLIP agents that set up, run, validate, "
              "and refine simulations on PNNL HPC from a natural-language goal")
    r = p.runs[0]; r.font.size = Pt(20); r.font.italic = True; r.font.color.rgb = MID

    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.45), Inches(12), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = ("Sarah Allec & Maxim Ziatdinov  ·  Physical & Computational "
              "Sciences Directorate, PNNL  ·  BES Autonomy Workshop")
    r = p.runs[0]; r.font.size = Pt(13); r.font.color.rgb = MUTED

    add_notes(slide, """
    This section covers the computational-simulation side of SciLink: how a
    natural-language research goal becomes a runnable, validated atomistic
    simulation on PNNL HPC (Deception, SLURM), and how the agents recover from
    their own mistakes. Recent work is foregrounded — the engine-neutral
    force-field path, pre-run validation gates, and an end-to-end aqueous
    electrolyte campaign.
    """)


# ── 2. the gap ────────────────────────────────────────────────────
def slide_gap(prs):
    slide = blank(prs)
    add_title(slide, "The gap SciLink closes")
    add_subtitle(slide,
                 "Turning a scientific question into a correct simulation on HPC "
                 "is expert-intensive, brittle, and slow to iterate.")
    add_bullets(slide, [
        (0, "Setting up a DFT / MD / MLIP calculation is specialist work — "
            "structure, force field, input deck, convergence settings, scheduler."),
        (0, "On HPC the failure cost is high: a one-character input error can run "
            "to completion and silently give the wrong physics, wasting node-hours."),
        (0, "Theory typically lags the experiment by weeks — the setup burden keeps "
            "simulation off the critical path of a discovery loop."),
        (0, "SciLink turns a natural-language goal into a runnable, VALIDATED "
            "simulation on HPC, and recovers from its own errors autonomously:"),
        (1, "picks the right method and engine for the science;"),
        (1, "authors and validates the input deck before submission;"),
        (1, "runs on PNNL HPC, then debugs failures and quality-checks results, "
            "re-running with corrected inputs until the physics holds."),
        (0, "Three autonomy levels — co-pilot / autopilot / autonomous — set who "
            "holds the acceptance gate: the human, partly, or the agent itself."),
    ])
    add_notes(slide, """
    The core value proposition for a simulation audience: SciLink collapses the
    expert setup burden and puts simulation inside the experimental loop instead
    of weeks behind it. The emphasis on "validated" is deliberate — on HPC the
    expensive failure mode is a run that completes but is physically wrong (see
    the VASP silent-typo story on the benchmark slide). Autonomy levels let a
    user dial trust from full oversight to hands-off.
    """)


# ── 3. architecture (adapted from build_overview_deck.slide_stack) ──
def slide_architecture(prs):
    slide = blank(prs)
    add_title(slide, "Engine-neutral simulation architecture")
    add_subtitle(slide,
                 "A router sends the goal to the right scale agent; each agent is "
                 "software-agnostic. Adding an engine is one skill bundle — no "
                 "core-code change.")

    add_node(slide, 2.1, 1.85, 7.9, 0.7,
             "Skill subsystem  ·  domain-keyed knowledge + engine bundles",
             sublabel="markdown (+ optional helpers) discovered and routed "
                      "automatically — domain experts extend without touching agents",
             fill=SKILL_FILL, border=SKILL_LINE, bold=True, font_size=13,
             sublabel_color=SKILL_LINE, sublabel_size=11)

    add_node(slide, 0.4, 4.0, 1.4, 0.9, "User goal",
             sublabel="natural language", fill=INPUT_FILL,
             border=RGBColor(0xC8, 0xA8, 0x3C), bold=True, font_size=12)

    add_node(slide, 2.1, 3.85, 3.1, 1.2, "Simulation Orchestrator",
             sublabel="plans · dispatches · adjudicates",
             fill=ORCH_FILL, border=ACCENT, bold=True, font_size=14)

    add_node(slide, 5.55, 4.05, 1.5, 0.85, "Router",
             sublabel="(scale, engine)", fill=ROUTER_FILL, border=ACCENT,
             bold=True, font_size=12)

    eng = 10
    # Four scale agents (periodic DFT · molecular QC · classical MD · MLIP), each
    # paired with the concrete input it generates. Rows are computed so the stack
    # stays clear of the captions below.
    agents = [
        ("Periodic-DFT agent", "VASP · QE · ABINIT · CP2K",
         "POSCAR · INCAR · KPOINTS", "(VASP example)"),
        ("Molecular-QC agent", "NWChem  (+ PySCF · ORCA in progress)",
         "job.nw deck", "(NWChem example)"),
        ("Classical-MD agent", "LAMMPS · GROMACS · OpenMM\nFF: OpenFF · AMBER",
         "MD deck + typed data file", "(LAMMPS example)"),
        ("MLIP agent", "MACE · CHGNet  (+ more in progress)",
         "Deployed potential", "(engine-neutral descriptor)"),
    ]
    a_top, a_h, a_gap, router_c = 2.95, 0.62, 0.13, 4.475
    for i, (name, engines, out, outsub) in enumerate(agents):
        ay = a_top + i * (a_h + a_gap)
        ac = ay + a_h / 2
        add_node(slide, 7.4, ay, 2.6, a_h, name, sublabel=engines,
                 fill=AGENT_FILL, border=ACCENT, bold=True, font_size=11,
                 sublabel_color=SKILL_LINE, sublabel_size=eng)
        add_node(slide, 10.25, ay, 2.6, a_h, out, sublabel=outsub,
                 fill=OUTPUT_FILL, border=MUTED, font_size=11, sublabel_size=9)
        add_arrow(slide, 7.05, router_c, 7.4, ac)     # router fans out
        add_arrow(slide, 10.0, ac, 10.25, ac)         # agent -> generated input

    add_arrow(slide, 3.65, 2.55, 3.65, 3.85, color=SKILL_LINE, dashed=True)
    add_arrow(slide, 8.7, 2.55, 8.7, 2.95, color=SKILL_LINE, dashed=True)
    add_arrow(slide, 1.8, 4.45, 2.1, 4.45)
    add_arrow(slide, 5.2, 4.45, 5.55, 4.475)

    add_caption(slide, 0.5, 6.2, 12.3, 0.4,
                "Solid arrows = data flow.    Dashed arrows = skill-bundle context "
                "feeding the router's decision and each agent's engine set.")
    add_caption(slide, 0.5, 6.55, 12.3, 0.4,
                "N + M, not N × M:  one new MLIP backend or force field is a single "
                "skill bundle, reused across every MD engine — not an integration "
                "per engine.", size=11)

    add_notes(slide, """
    This is the "foundation-agent" shape: one agent per scale (periodic DFT,
    molecular quantum chemistry, classical MD, MLIP-driven MD), each specialized
    at runtime to a specific engine by a pluggable skill bundle. Molecular QC is
    the newest agent -- NWChem is fully wired (it emits a single job.nw deck),
    with PySCF/ORCA in progress. The engine-neutral descriptors are the
    load-bearing idea — e.g. the MLIP agent emits a "deployed potential"
    (backend, model, elements, calculator spec) that the MD agent consumes
    without importing any MLIP code. Recent work: the classical-MD force-field
    path is now engine-neutral too (OpenFF Interchange core produces runnable
    LAMMPS / GROMACS / OpenMM inputs from one parameterized system). MACE and
    CHGNet are the shipped MLIP backends; additional backends (UMA, Orb, DeePMD)
    are in review.
    """)


# ── 4. self-refinement + pre-run gates ────────────────────────────
def slide_refinement(prs):
    slide = blank(prs)
    add_title(slide, "Self-refinement + pre-run validation")
    add_subtitle(slide,
                 "Generate runs once; the loop is around Run. A cheap pre-run gate "
                 "catches deck errors before the expensive HPC job; debuggers and "
                 "quality checks close the loop after.")

    top_y, top_h = 2.15, 0.85
    mid_y = top_y + top_h / 2

    add_node(slide, 0.5, top_y, 1.9, top_h, "Generate",
             sublabel="structure + inputs", fill=ORCH_FILL, border=ACCENT,
             bold=True, font_size=13, sublabel_size=9)
    add_node(slide, 2.7, top_y, 2.5, top_h, "Pre-run gate",
             sublabel="setup + observable coverage\n(cheap dry-run twin)",
             fill=FILL_ALT, border=SKILL_LINE, bold=True, font_size=13,
             sublabel_color=SKILL_LINE, sublabel_size=8)
    add_node(slide, 5.5, top_y, 2.7, top_h, "Run on HPC",
             sublabel="optim → equilib → production", fill=ORCH_FILL,
             border=ACCENT, bold=True, font_size=13, sublabel_size=9)
    add_node(slide, 8.5, top_y, 2.5, top_h, "Quality check",
             sublabel="fires after every phase", fill=ORCH_FILL, border=ACCENT,
             bold=True, font_size=13, sublabel_size=9)
    add_node(slide, 11.3, top_y, 1.5, top_h, "Done",
             sublabel="(all phases pass)", fill=OUTPUT_FILL, border=MUTED,
             bold=True, font_size=13, sublabel_size=9)

    add_arrow(slide, 2.4, mid_y, 2.7, mid_y)
    add_arrow(slide, 5.2, mid_y, 5.5, mid_y)
    add_arrow(slide, 8.2, mid_y, 8.5, mid_y)
    add_arrow(slide, 11.0, mid_y, 11.3, mid_y)
    add_caption(slide, 8.2, mid_y - 0.30, 0.55, 0.22, "no error",
                size=9, color=MID)
    add_caption(slide, 11.0, mid_y - 0.30, 0.55, 0.22, "all pass",
                size=9, color=MID)

    # pre-run gate self-loop (fix before the HPC job)
    add_arrow(slide, 3.95, top_y + top_h, 3.95, top_y + top_h + 0.45,
              color=SKILL_LINE, dashed=True, arrow=False)
    add_arrow(slide, 3.95, top_y + top_h + 0.45, 3.35, top_y + top_h + 0.45,
              color=SKILL_LINE, dashed=True, arrow=False)
    add_arrow(slide, 3.35, top_y + top_h + 0.45, 3.35, top_y + top_h,
              color=SKILL_LINE, dashed=True, arrow=True)
    add_caption(slide, 2.6, top_y + top_h + 0.48, 2.6, 0.22,
                "fix deck, re-check (seconds)", size=9, color=SKILL_LINE)

    # debugger box (engine error branch)
    dbg_y, dbg_h = 4.35, 0.75
    add_node(slide, 5.6, dbg_y, 2.5, dbg_h, "Debugger",
             sublabel="parses engine error → corrected inputs",
             fill=ORCH_FILL, border=ACCENT, bold=True, font_size=13,
             sublabel_size=8)
    add_arrow(slide, 7.1, top_y + top_h, 7.1, dbg_y)
    add_caption(slide, 7.2, 3.35, 1.3, 0.22, "engine error", size=9, color=MID)
    add_arrow(slide, 6.0, dbg_y, 6.0, top_y + top_h,
              color=SKILL_LINE, dashed=True, arrow=True)
    add_caption(slide, 4.3, 3.35, 1.6, 0.22, "corrected inputs",
                size=9, color=SKILL_LINE, align=PP_ALIGN.RIGHT)

    # quality -> run (refine if questionable), routed above the debugger
    qx, loop_y, refine_x = 9.75, 4.05, 5.75
    add_arrow(slide, qx, top_y + top_h, qx, loop_y,
              color=SKILL_LINE, dashed=True, arrow=False)
    add_arrow(slide, qx, loop_y, refine_x, loop_y,
              color=SKILL_LINE, dashed=True, arrow=False)
    add_arrow(slide, refine_x, loop_y, refine_x, top_y + top_h,
              color=SKILL_LINE, dashed=True, arrow=True)
    add_caption(slide, 6.7, loop_y + 0.03, 3.6, 0.22,
                "refine if questionable", size=9, color=SKILL_LINE)

    add_bullets(slide, [
        (0, "Pre-run gate (recent work): a seconds-long dry-run twin checks that "
            "the deck starts cleanly AND emits every output the research goal "
            "needs — catching silent setup errors before a 12-hour HPC job."),
        (0, "Same loop across DFT, classical MD, and MLIP-driven MD — only the "
            "agents in each slot differ; the iteration is engine-agnostic."),
    ], top=Inches(5.45), size=13)

    add_notes(slide, """
    The self-refinement skeleton: Generate -> Run -> branch. On an engine error a
    debugger parses the log and proposes corrected inputs; on a completed-but-
    questionable phase a quality check triggers a refine. Both loop back to Run.

    The recent addition is the pre-run gate. It runs a cheap "dry-run twin" of
    the deck (engine setup only, seconds on the node) and fixes the real deck
    before the expensive job. Two dimensions: (1) does setup start cleanly, and
    (2) does the deck emit the observables the goal actually needs. The second is
    new and directly motivated by the electrolyte campaign — some decks ran
    fine but never logged the stress needed for viscosity, so the property was
    unrecoverable after a 12-hour run. The gate now catches that up front.
    """)


# ── 5. flagship application: aqueous Zn electrolyte on HPC ─────────
def slide_application(prs):
    slide = blank(prs)
    add_title(slide, "End-to-end on PNNL HPC:  aqueous Zn-electrolyte campaign")
    add_subtitle(slide,
                 "One high-level goal per system → full setup, run, and refine on "
                 "Deception (SLURM). Validating against NMR / transport measurements "
                 "of Zn(OTf)₂ in water / sulfone mixtures.")

    # pipeline strip
    stages = [
        ("Goal", "\"...run MD of 1 M\nZn(OTf)₂ in water\"", INPUT_FILL,
         RGBColor(0xC8, 0xA8, 0x3C)),
        ("Structure", "Packmol box\n+ components", AGENT_FILL, ACCENT),
        ("Force field", "OpenFF Interchange\nNAGL charges · ion vdW", AGENT_FILL, ACCENT),
        ("LAMMPS deck", "typed data file\n+ MD script", AGENT_FILL, ACCENT),
        ("Run + refine", "Deception (SLURM)\nmulti-phase MD", ORCH_FILL, ACCENT),
        ("Properties", "density · D(Zn)\nviscosity · RDF", OUTPUT_FILL, MUTED),
    ]
    x = 0.5
    w = 1.95
    gap = 0.12
    y = 1.95
    for i, (lab, sub, fill, border) in enumerate(stages):
        add_node(slide, x, y, w, 0.95, lab, sublabel=sub, fill=fill, border=border,
                 bold=True, font_size=12, sublabel_size=8)
        if i < len(stages) - 1:
            add_arrow(slide, x + w, y + 0.475, x + w + gap, y + 0.475)
        x += w + gap

    # results table
    add_caption(slide, 0.5, 3.15, 12.3, 0.26,
                "Selected results  (298.15 K; computed vs. experiment)",
                size=12, italic=False, color=ACCENT, align=PP_ALIGN.LEFT)
    tbl = slide.shapes.add_table(
        4, 4, Inches(0.5), Inches(3.45), Inches(8.3), Inches(1.9)).table
    for c, wd in enumerate([3.5, 1.7, 1.6, 1.5]):
        tbl.columns[c].width = Inches(wd)
    fill_table(tbl, [
        ["property", "computed", "experiment", "verdict"],
        ["Mass density (g/cm³), S1→S4", "1.16 → 1.19", "increasing", "trend ✓"],
        ["Zn²⁺ self-diffusion vs. sulfone frac.", "decreasing", "decreasing", "trend ✓"],
        ["Shear viscosity (cP), S3 / S4", "0.8 / 1.6", "3.4 / 4.5", "3–4× low"],
    ])
    style_table(tbl, font_size=11)

    # honest-framing callout
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(3.45), Inches(3.8), Inches(1.9))
    bg.fill.solid(); bg.fill.fore_color.rgb = FILL_ALT
    bg.line.color.rgb = AMBER; bg.shadow.inherit = False
    add_caption(slide, 9.15, 3.52, 3.5, 0.3, "What this shows",
                size=12, italic=False, color=AMBER, align=PP_ALIGN.LEFT)
    add_caption(slide, 9.15, 3.85, 3.5, 1.45,
                "The full pipeline runs autonomously on HPC and reproduces the "
                "right physical trends. Absolute transport is force-field-limited "
                "(TIP3P water is ~3× too fluid) — a known model limit, not an "
                "agent failure. Better water models are a drop-in skill swap.",
                size=10, italic=False, color=DARK, align=PP_ALIGN.LEFT)

    add_bullets(slide, [
        (0, "Live driver of the pre-run gate: some systems ran to completion but "
            "never logged the stress needed for viscosity — exactly the "
            "observable-coverage gap the gate now catches before submission."),
    ], top=Inches(5.5), size=12)

    add_notes(slide, """
    This is the flagship recent-work application. Five systems (S1-S5): 1 M
    Zn(OTf)2 in water / ethyl-isopropyl-sulfone mixtures, 298.15 K, for
    comparison with NMR relaxation and viscosity data. Each is driven by a single
    high-level request; the agent derives chemistry, box, force field, and deck.

    Grounded numbers: computed mass densities 1.159 (S1) -> 1.188 (S4) g/cm3,
    increasing with sulfone fraction (physically correct). Zn self-diffusion
    decreases with sulfone fraction (correct trend). Green-Kubo shear viscosity
    ~0.8 cP (S3) / ~1.6 cP (S4) vs experimental ~3.4 / ~4.5 cP -- 3-4x low,
    consistent with TIP3P water being ~3x too fluid, compounded by the density
    offset. The S1->S5 increasing trend is reproduced.

    Honest framing for a BES audience: the value is the autonomous end-to-end
    HPC pipeline and correct trends; absolute transport accuracy is bounded by
    the classical force field, and improving it is a skill swap (a better water
    model), not new agent code.

    Tie-back to slide 4: S3/S4 emitted the stress tensor (viscosity computable);
    S2/S5 did not (no stress.dat) -- the concrete failure that motivated the
    observable-coverage dimension of the pre-run gate.
    """)


# ── 6. benchmarks + evaluation methodology ────────────────────────
def slide_benchmarks(prs):
    slide = blank(prs)
    add_title(slide, "Validation + rigorous evaluation")
    add_subtitle(slide,
                 "Can the agent pick the right method, write a correct deck, and "
                 "run it? Benchmarked on PNNL-relevant systems on Deception — "
                 "and stress-tested for consistency, not just single-shot accuracy.")

    stat_card(slide, 0.5, "100 %",
              "router joint (scale + engine) accuracy\n21 scorable prompts, 3 models")
    stat_card(slide, 4.7, "8 / 9",
              "DFT cells validated\n6 crystals ≤ 1.3 % of expt lattice const.")
    stat_card(slide, 8.9, "5 / 6",
              "MLIP→MD trajectories complete\n(MACE, engine-neutral handoff)")

    # methodology callout
    add_caption(slide, 0.5, 3.0, 12.3, 0.28,
                "Single-shot accuracy is necessary but not sufficient — we "
                "N-shot sample every decision",
                size=12, italic=False, color=ACCENT, align=PP_ALIGN.LEFT)

    tbl = slide.shapes.add_table(
        4, 4, Inches(0.5), Inches(3.35), Inches(7.7), Inches(2.0)).table
    for c, wd in enumerate([3.2, 1.5, 1.5, 1.5]):
        tbl.columns[c].width = Inches(wd)
    fill_table(tbl, [
        ["input-deck reliability (8 trials × 3 prompts)", "opus-4-6", "sonnet-4-5", ""],
        ["INCAR typo rate", "0.21", "0.46", "lower = better"],
        ["physics-setting stability", "0.96", "0.88", "higher = better"],
        ["router pick stability", "1.00", "1.00", "self-consistent"],
    ])
    style_table(tbl, font_size=11)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.4), Inches(3.35), Inches(4.4), Inches(2.0))
    bg.fill.solid(); bg.fill.fore_color.rgb = FILL_ALT
    bg.line.color.rgb = SKILL_LINE; bg.shadow.inherit = False
    add_caption(slide, 8.55, 3.42, 4.1, 0.3, "Why it matters",
                size=12, italic=False, color=SKILL_LINE, align=PP_ALIGN.LEFT)
    add_caption(slide, 8.55, 3.75, 4.1, 1.55,
                "VASP silently ignores an unknown INCAR tag — a one-letter typo "
                "(ISPN for ISPIN) disables spin polarization with no error, and Fe "
                "comes back non-magnetic. The run 'succeeds' but the physics is "
                "wrong. Model choice measurably changes this risk, so a pre-submit "
                "syntax check now gates every deck.",
                size=10, italic=False, color=DARK, align=PP_ALIGN.LEFT)

    add_caption(slide, 0.5, 5.55, 12.3, 0.5,
                "Six bulk crystals (Cu · Si · C · MgO · LiCoO₂ · UO₂) and "
                "magnetic Fe all within 1.3 % of experiment; Pt(111)+CO adsorption "
                "within 3 meV. Router: 100 % joint accuracy across opus-4-6, "
                "opus-4-7, and sonnet-4-5 on 21 scorable prompts (3 gas-phase "
                "quantum-chemistry prompts excluded as a known capability gap).",
                size=11)

    add_notes(slide, """
    Three benchmark suites, all on Deception, results in benchmark/outputs/:

    Router (test_router): 24 PNNL-relevant prompts -> top-1 (scale, engine).
    100% joint accuracy on the 21 scorable prompts across three models
    (opus-4-6, opus-4-7, sonnet-4-5). Three molecular-DFT prompts excluded
    (no molecular-DFT agent yet -- an honest capability gap).

    DFT cells (test_dft): 8/9 pass. Six bulk crystals within <=1.3% of the
    experimental lattice constant (Cu 0.39, Si 0.69, C 0.18, MgO 0.83,
    LiCoO2 1.30, UO2 0.59; Fe -1.24 after the pre-submit validator). Pt(111)+CO
    adsorption energy within 3 meV. One failure: TiO2(101)+H2O slab.

    MLIP->MD (test_mlip): 5 of 6 solid-state cells deployed MACE and completed a
    1000-step NVT trajectory; one liquid box skipped (Packmol materialization
    gap). First end-to-end MLIP->MD on the stack via the engine-neutral
    deployed-potential handoff.

    Evaluation methodology (the BES-relevant point): single-shot accuracy hides
    variance. We N-shot sample (8 trials x 3 prompts x 2 models). Input-deck
    reliability is measurably model-dependent -- INCAR typo rate 0.21 (opus-4-6)
    vs 0.46 (sonnet-4-5); physics-setting stability 0.96 vs 0.88. Router picks
    are perfectly self-consistent for both. The lesson: rigorous autonomy
    evaluation means measuring consistency, and pairing it with deterministic
    pre-submit validation so a model's off-days can't reach the cluster.

    The Fe/ISPIN story is the canonical "converged but wrong" failure and the
    single best motivation for pre-run validation to a simulation audience.
    """)


# ── 7. successes / challenges / lessons ───────────────────────────
def slide_lessons(prs):
    slide = blank(prs)
    add_title(slide, "Successes, challenges & lessons learned")
    add_subtitle(slide,
                 "Where autonomous atomistic simulation on HPC stands today — "
                 "and what the work has taught us.")

    col_w = 4.0
    xs = [0.5, 4.65, 8.8]
    heads = [("Successes", GREEN), ("Challenges", AMBER), ("Lessons", ACCENT)]
    body = [
        [  # successes
            "End-to-end autonomous runs on PNNL HPC (Deception, SLURM) across "
            "DFT, classical MD, and MLIP-driven MD.",
            "100 % method+engine routing; DFT within 1.3 % of experiment; first "
            "engine-neutral MLIP→MD handoff.",
            "Engine-neutral force fields + skills: a new engine or FF is one "
            "bundle, reused across engines.",
        ],
        [  # challenges
            "Classical force fields bound absolute accuracy (e.g. transport with "
            "TIP3P water) — physics, not agent, limits.",
            "Capability gaps remain (e.g. gas-phase molecular DFT; some liquid-box "
            "materialization).",
            "Reliability is model-dependent — deck-error rates differ measurably "
            "between models.",
        ],
        [  # lessons
            "Validate before you spend: a seconds-long pre-run gate saves "
            "12-hour HPC jobs from silent errors.",
            "Evaluate consistency, not just single-shot accuracy — N-shot "
            "sampling surfaces what one run hides.",
            "Keep engine specifics in skills, behind neutral contracts — that is "
            "what makes the platform extensible.",
        ],
    ]
    for (head, color), x, items in zip(heads, xs, body):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.75), Inches(col_w), Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background(); bar.shadow.inherit = False
        p = bar.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        p.text = head
        r = p.runs[0]; r.font.size = Pt(16); r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        add_bullets(slide, [(0, it) for it in items],
                    left=Inches(x), top=Inches(2.5), width=Inches(col_w),
                    height=Inches(4.0), size=13)

    add_caption(slide, 0.5, 6.7, 12.3, 0.4,
                "SciLink · PNNL Physical & Computational Sciences Directorate · "
                "PCSD Seed LDRD + Foundational Autonomy Investment",
                size=10)

    add_notes(slide, """
    Maps directly to the colleague's "successes, challenges, lessons" topic, and
    threads the integration and cross-facility themes.

    Successes: the platform actually runs end-to-end on PNNL HPC across three
    simulation scales, with strong validation numbers and a genuinely
    engine-neutral extensibility story.

    Challenges (stated honestly): classical FFs cap absolute accuracy; real
    capability gaps remain; and model choice affects reliability.

    Lessons: (1) cheap pre-run validation before expensive HPC jobs is the single
    highest-leverage robustness investment; (2) rigorous autonomy evaluation
    means measuring consistency, not just single-shot accuracy; (3) the
    engine-neutral-skills architecture is what keeps all of this extensible
    without rewriting agents.
    """)


def main():
    prs = Presentation()
    set_widescreen(prs)
    slide_title(prs)
    slide_gap(prs)
    slide_architecture(prs)
    slide_refinement(prs)
    slide_application(prs)
    slide_benchmarks(prs)
    slide_lessons(prs)
    out = "scilink_bes_sim_deck.pptx"
    prs.save(out)
    print(f"wrote {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
